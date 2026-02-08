"""
Generate Executive Summary PowerPoint v2.0 for ASPR Photo Repository.
14-slide post-deployment briefing with ASPR + Leidos branding.

Run:  python scripts/generate_exec_summary_pptx.py
Requires: pip install python-pptx
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "ASPR_Photo_Repository_Executive_Summary.pptx"

# ── Brand Colors ──────────────────────────────────────────────────────
BLUE_DARK     = RGBColor(0x06, 0x2E, 0x61)
BLUE_PRIMARY  = RGBColor(0x15, 0x51, 0x97)
BLUE_MEDIUM   = RGBColor(0x24, 0x77, 0xBD)
GOLD          = RGBColor(0xAA, 0x64, 0x04)
RED           = RGBColor(0x99, 0x00, 0x00)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY     = RGBColor(0x32, 0x32, 0x32)
GOLD_LIGHT    = RGBColor(0xFB, 0xD0, 0x98)

ROW_EVEN      = RGBColor(0x0A, 0x3D, 0x7A)
ROW_ODD       = RGBColor(0x08, 0x35, 0x6E)
CARD_BG       = RGBColor(0x0A, 0x3D, 0x7A)
CARD_BORDER   = RGBColor(0x15, 0x51, 0x97)
MUTED         = RGBColor(0x99, 0x99, 0x99)

# ── Logo Paths ────────────────────────────────────────────────────────
ASPR_LOGO = ROOT / "public" / "aspr-logo-blue.png"
HHS_LOGO = ROOT / "public" / "hhs_longlogo_white.png"
LEIDOS_LOGO = Path(
    r"C:\Users\ravinder.mathaudhu\OneDrive - HHS Office of the Secretary"
    r"\Documents\Projects\New folder\Leidos-Logo-Suite\Leidos-Logo-Suite"
    r"\02-Digital\03-Raster-PNG\Leidos-logo-horz-full-rgb-@2x.png"
)


# ══════════════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════════════

def add_dark_bg(slide, color=BLUE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.06), color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), top, Inches(13.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, text="HHS/ASPR \u2014 For Official Use Only | Leidos"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0),
                                      Inches(12.333), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER


def add_title_text(slide, text, left, top, width, height,
                   font_size=Pt(36), color=WHITE, bold=True,
                   alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_slide_header(slide, title):
    add_title_text(slide, title,
                   Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
                   font_size=Pt(32), color=WHITE, bold=True)
    add_accent_bar(slide, top=Inches(1.4), height=Inches(0.04), color=GOLD)


def add_bullet_slide(slide, title, bullets, title_color=WHITE,
                     bullet_color=WHITE, font_size=Pt(18)):
    add_slide_header(slide, title)

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8),
                                      Inches(11), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.font.size = font_size
        p.font.color.rgb = bullet_color
        p.space_after = Pt(12)
        p.level = 0

    return tf


def add_table_slide(slide, title, headers, rows, col_widths=None,
                    font_hdr=Pt(14), font_row=Pt(13)):
    add_slide_header(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), Inches(1.8),
        Inches(11.5), Inches(0.4) * n_rows
    )
    table = table_shape.table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(Inches(11.5) * w / total)

    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.size = font_hdr
            p.font.color.rgb = WHITE
            p.font.bold = True

    for ri, row_data in enumerate(rows):
        bg = ROW_EVEN if ri % 2 == 0 else ROW_ODD
        for ci, text in enumerate(row_data[:n_cols]):
            cell = table.cell(ri + 1, ci)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = font_row
                p.font.color.rgb = WHITE

    return table


def add_kpi_cards(slide, cards):
    """Row of metric cards. Each card = (number, label, sublabel)."""
    n = len(cards)
    card_w = Inches(2.2)
    gap = Inches(0.25)
    total_w = n * card_w + (n - 1) * gap
    start_x = (Inches(13.333) - total_w) / 2
    y = Inches(2.2)

    for i, (number, label, sublabel) in enumerate(cards):
        x = start_x + i * (card_w + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(2.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.5)

        add_title_text(slide, number,
                       x + Inches(0.15), y + Inches(0.3),
                       card_w - Inches(0.3), Inches(0.9),
                       font_size=Pt(44), color=GOLD_LIGHT, bold=True,
                       alignment=PP_ALIGN.CENTER)
        add_title_text(slide, label,
                       x + Inches(0.15), y + Inches(1.2),
                       card_w - Inches(0.3), Inches(0.5),
                       font_size=Pt(16), color=WHITE, bold=True,
                       alignment=PP_ALIGN.CENTER)
        add_title_text(slide, sublabel,
                       x + Inches(0.15), y + Inches(1.75),
                       card_w - Inches(0.3), Inches(0.8),
                       font_size=Pt(11), color=MUTED, bold=False,
                       alignment=PP_ALIGN.CENTER)


def add_two_col_features(slide, title, left_title, left_items,
                         right_title, right_items):
    """Two-column feature list with gold sub-headers."""
    add_slide_header(slide, title)

    for col_idx, (col_title, items) in enumerate([
        (left_title, left_items), (right_title, right_items)
    ]):
        x = Inches(0.8) if col_idx == 0 else Inches(7.0)
        add_title_text(slide, col_title,
                       x, Inches(1.7), Inches(5.5), Inches(0.4),
                       font_size=Pt(20), color=GOLD_LIGHT, bold=True)

        txBox = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.2),
                                          Inches(5.4), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"\u2022  {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)


# ══════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION  (14 slides)
# ══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)   # Widescreen 16:9
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


# ── SLIDE 1: TITLE ──────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_accent_bar(slide, top=Inches(0), height=Inches(0.08), color=GOLD)

if ASPR_LOGO.exists():
    slide.shapes.add_picture(str(ASPR_LOGO), Inches(0.8), Inches(0.4),
                             height=Inches(1.0))
if LEIDOS_LOGO.exists():
    slide.shapes.add_picture(str(LEIDOS_LOGO), Inches(10.5), Inches(0.4),
                             height=Inches(0.7))

add_title_text(slide, "Executive Summary",
               Inches(0.8), Inches(2.2), Inches(11), Inches(1.0),
               font_size=Pt(48), color=WHITE, bold=True)
add_title_text(slide, "ASPR Photo Repository Application",
               Inches(0.8), Inches(3.2), Inches(11), Inches(0.7),
               font_size=Pt(28), color=GOLD_LIGHT, bold=False)

add_accent_bar(slide, top=Inches(4.1), height=Inches(0.04), color=GOLD)

metadata_lines = [
    "U.S. Department of Health and Human Services",
    "Administration for Strategic Preparedness and Response (ASPR)",
    "",
    "Prepared by: HHS ASPR / Leidos",
    "Date: February 7, 2026  |  Version 2.0",
    "Status: DEPLOYED TO PRODUCTION",
    "Classification: For Official Use Only (FOUO)",
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5),
                                  Inches(11), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(metadata_lines):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC) if line else WHITE
    if "Department" in line or "Administration" in line:
        p.font.color.rgb = WHITE
        p.font.size = Pt(18)
    if "DEPLOYED" in line:
        p.font.color.rgb = GOLD_LIGHT
        p.font.bold = True

add_footer(slide)


# ── SLIDE 2: PURPOSE & MISSION ──────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_bullet_slide(slide, "Purpose & Mission", [
    "Enable ASPR field teams to securely capture, upload, and manage disaster-related "
    "photographs during incident response operations",
    "Provide rapid photo documentation capability deployable within hours of incident "
    "activation \u2014 now live in production with full CDN acceleration",
    "Replace ad-hoc photo collection methods (email, shared drives, USB) with a "
    "purpose-built, secure web application accessible via PIN, Entra ID SSO, "
    "Login.gov, and ID.me",
    "Support incident accountability with geotagged, timestamped, EXIF-enriched "
    "photographic evidence and full admin audit trail",
    "Operate within the HHS/ASPR security boundary with Azure Front Door WAF "
    "(OWASP 3.2), Private Link network isolation, and NIST SP 800-53 alignment",
])
add_footer(slide)


# ── SLIDE 3: PLATFORM HIGHLIGHTS (KPI CARDS) ────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_slide_header(slide, "What We Built \u2014 Platform Highlights")

add_kpi_cards(slide, [
    ("17+", "API Endpoints", "REST API with full\nCRUD + bulk operations"),
    ("4", "Auth Methods", "PIN, Entra ID SSO,\nLogin.gov, ID.me"),
    ("3", "Image Renditions", "thumb_sm, thumb_md,\nweb (all WebP)"),
    ("8", "Database Tables", "SQL + audit log\n+ EXIF + tags"),
    ("10+", "Admin Components", "Photo grid, editor,\ntags, bulk ops"),
])

add_footer(slide)


# ── SLIDE 4: KEY CAPABILITIES ───────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_slide_header(slide, "Key Capabilities")

capabilities_left = [
    ("Multi-Auth Security", "PIN + JWT (field), Entra ID SSO (admin),\n"
     "Login.gov & ID.me (external), rate limiting"),
    ("Photo Upload Wizard", "6-step guided upload with animated progress,\n"
     "GPS capture, incident tagging, batch support"),
    ("Admin Photo Grid", "Virtualized grid with search, filters,\n"
     "status badges, bulk select, cursor pagination"),
    ("Photo Editor", "Crop (aspect presets), rotate 90\u00b0,\n"
     "flip H/V, rendition regeneration"),
]
capabilities_right = [
    ("Tag System", "Categorized tags (status, priority, type,\n"
     "timeline, custom) with autocomplete"),
    ("EXIF Extraction", "Camera make/model, lens, aperture, ISO,\n"
     "shutter speed, GPS altitude, date taken"),
    ("Bulk Operations", "Multi-select delete, tag assignment,\n"
     "status change, ZIP download"),
    ("Session Management", "Create/revoke PINs, view photo counts,\n"
     "storage usage, team tracking"),
]

for col_idx, caps in enumerate([capabilities_left, capabilities_right]):
    x = Inches(0.8) if col_idx == 0 else Inches(7.0)
    for i, (cap_title, desc) in enumerate(caps):
        y = Inches(1.9) + Inches(1.25) * i
        add_title_text(slide, cap_title,
                       x, y, Inches(5.5), Inches(0.4),
                       font_size=Pt(18), color=GOLD_LIGHT, bold=True)
        add_title_text(slide, desc,
                       x, y + Inches(0.38), Inches(5.5), Inches(0.75),
                       font_size=Pt(14), color=WHITE, bold=False)

add_footer(slide)


# ── SLIDE 5: ARCHITECTURE OVERVIEW ──────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_table_slide(slide, "Architecture Overview",
    ["Layer", "Component", "Technology", "Purpose"],
    [
        ["Application", "Web Framework", "Next.js 16.1.6 (React 19)", "Full-stack SSR + API routes"],
        ["Application", "UI / Design", "Tailwind CSS 4 + shadcn/ui", "Glassmorphic component system"],
        ["Application", "Image Pipeline", "Sharp 0.34 + exifr", "Multi-rendition WebP + EXIF"],
        ["Security", "WAF", "Azure Front Door WAF", "OWASP DRS 2.1 + Bot Protection"],
        ["Security", "Authentication", "Auth.js v5 + bcrypt + JWT", "Multi-provider auth system"],
        ["Network", "CDN", "Azure Front Door Premium", "Global edge caching + SSL"],
        ["Network", "Private Link", "Azure Private Endpoints", "VNet isolation (blob + app)"],
        ["Data", "Database", "Azure SQL Server", "Sessions, photos, tags, audit"],
        ["Data", "Blob Storage", "Azure Blob Storage", "Photo originals + renditions"],
        ["Data", "Key Vault", "Azure Key Vault", "Secrets management"],
        ["Hosting", "App Service", "Linux / Node.js 22", "Standalone Next.js runtime"],
        ["CI/CD", "Pipeline", "GitHub Actions", "ZipDeploy + post-deploy migrate"],
    ],
    col_widths=[13, 18, 30, 39],
    font_hdr=Pt(13), font_row=Pt(12),
)
add_footer(slide)


# ── SLIDE 6: SECURITY POSTURE ───────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_bullet_slide(slide, "Security Posture", [
    "FIPS 199 MODERATE categorization \u2014 appropriate for operational "
    "incident photography",
    "Azure Front Door WAF (OWASP DRS 2.1 + Microsoft Bot Manager) in "
    "Prevention mode protecting all application traffic",
    "Network isolation via Private Endpoints \u2014 Blob Storage, SQL, "
    "and Key Vault on VNet; App Service behind Private Link origins",
    "OWASP Top 10 (2021) fully addressed \u2014 injection prevention, "
    "access control, cryptographic protections, security misconfiguration",
    "NIST SP 800-63B compliant PIN generation (CSPRNG) with bcrypt "
    "storage (10 salt rounds)",
    "Comprehensive rate limiting \u2014 5 PIN attempts/min (15-min lockout), "
    "3 admin attempts (30-min lockout), 50 uploads/hour",
    "Hardened HTTP headers \u2014 HSTS, CSP, X-Frame-Options, "
    "Permissions-Policy on all routes",
    "Immutable admin audit log \u2014 all operations recorded with entity, "
    "performer email, IP address, timestamp",
    "Signed image URLs (HMAC-SHA256) \u2014 24-hour expiry, no JWT "
    "exposure in query strings",
], font_size=Pt(16))
add_footer(slide)


# ── SLIDE 7: ADMIN DASHBOARD SHOWCASE ───────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_two_col_features(slide, "Admin Dashboard \u2014 Full Photo Management",
    "Management Features", [
        "Photo grid with virtual scrolling (100/page cursor pagination)",
        "Search by filename, filter by incident/status/date/session/tags",
        "Photo detail sidebar with inline metadata editing",
        "Photo editor: crop with aspect presets, rotate, flip",
        "Rendition auto-regeneration after edits (thumb_sm, thumb_md, web)",
        "Admin bulk upload panel (drag-and-drop, up to 50 files)",
        "Dashboard statistics: totals, incidents, daily volume, top teams",
    ],
    "Organization & Operations", [
        "Tag system: status, priority, type, timeline, custom categories",
        "Tag autocomplete with category filtering and color coding",
        "Bulk operations: delete, tag assign/remove, status change",
        "Bulk download: client-side ZIP via signed URLs",
        "EXIF data: camera make/model, lens, aperture, ISO, GPS, date",
        "Session manager: create/revoke PINs, usage stats per team",
        "Audit log: entity type, action, performer, IP, details JSON",
    ],
)
add_footer(slide)


# ── SLIDE 8: CDN & PERFORMANCE ──────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_table_slide(slide, "CDN & Performance Architecture",
    ["Component", "Configuration", "Details"],
    [
        ["Front Door Profile", "Premium_AzureFrontDoor", "cdn-ociomicro-premium-eus2-01 (shared)"],
        ["App Endpoint", "cdn-asprphotos-app", "All app routes (/*), HTTPS-only"],
        ["Blob Endpoint", "cdn-asprphotos", "Rendition images (/renditions/*), HTTPS-only"],
        ["WAF Policy", "wafAsprPhotos", "OWASP DRS 2.1 + Bot Protection, Prevention mode"],
        ["App Origin", "Private Link", "App Service via approved Private Endpoint"],
        ["Blob Origin", "Private Link", "Blob Storage via approved Private Endpoint"],
        ["Health Probe", "/api/health", "Every 30s \u2014 HTTP 200 + JSON status check"],
        ["Image Renditions", "3 variants/photo", "thumb_sm 200x150, thumb_md 400x300, web 1200px"],
        ["Cache Strategy", "7-day immutable", "Static assets + hero images; API routes no-cache"],
    ],
    col_widths=[22, 28, 50],
    font_row=Pt(12),
)
add_footer(slide)


# ── SLIDE 9: CI/CD PIPELINE ────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_bullet_slide(slide, "CI/CD Pipeline \u2014 Automated Deployment", [
    "1.  Trigger: Push to main branch or manual workflow_dispatch",
    "2.  Build: Node.js 22.x \u2014 npm install + npm run build "
    "(Next.js standalone output)",
    "3.  Package: Copy .next/static + public/ into .next/standalone artifact",
    "4.  Deploy: azure/webapps-deploy@v2 via publish profile "
    "(ZipDeploy to SCM endpoint)",
    "5.  Target: app-aspr-photos in rg-ocio-microsites-eus2-01",
    "6.  Post-Deploy: POST /api/admin/migrate with x-admin-token "
    "for database schema migrations",
    "7.  Health: /api/health endpoint polled every 30s by "
    "Front Door health probe",
    "8.  Runtime: node server.js (configured on App Service, "
    "not in workflow)",
    "9.  Secrets: AZURE_WEBAPP_PUBLISH_PROFILE stored as "
    "GitHub Actions encrypted secret",
], font_size=Pt(16))
add_footer(slide)


# ── SLIDE 10: TIMELINE & MILESTONES ─────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_table_slide(slide, "Timeline & Milestones",
    ["Phase", "Timeline", "Status", "Key Deliverables"],
    [
        ["1. Requirements & Design", "Jan 2026", "COMPLETE",
         "SRS v2.0, SDD, Security Plan, architecture review"],
        ["2. Core Development", "Jan\u2013Feb 2026", "COMPLETE",
         "DB schema, PIN auth, upload API, gallery, wizard"],
        ["3. Security Hardening", "Feb 2026", "COMPLETE",
         "bcrypt, JWT, rate limiting, signed URLs, CSP headers"],
        ["4. Admin Dashboard", "Feb 2026", "COMPLETE",
         "Photo grid, editor, bulk ops, tags, EXIF, sessions"],
        ["5. Infrastructure & CDN", "Feb 2026", "COMPLETE",
         "Front Door Premium, WAF, Private Link, CDN endpoints"],
        ["6. CI/CD & Deployment", "Feb 2026", "COMPLETE",
         "GitHub Actions, ZipDeploy, post-deploy migrate"],
        ["7. UI/UX Polish", "Feb 2026", "COMPLETE",
         "Glassmorphic design, animations, preloader, transitions"],
        ["8. Documentation", "Feb 2026", "COMPLETE",
         "6-document suite + PPTX + Project Plan XML"],
        ["9. UAT & ATO", "Feb\u2013Mar 2026", "IN PROGRESS",
         "User acceptance testing, security review, ATO package"],
        ["10. Production Ops", "Mar 2026+", "PLANNED",
         "Monitoring, training, field pilot, v1.1 planning"],
    ],
    col_widths=[22, 13, 12, 53],
    font_hdr=Pt(13), font_row=Pt(12),
)
add_footer(slide)


# ── SLIDE 11: DOCUMENT PACKAGE ──────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_table_slide(slide, "Professional Document Package",
    ["#", "Document", "Version", "Description"],
    [
        ["01", "Software Requirements Specification", "v2.0",
         "Functional & non-functional requirements, data model, API spec"],
        ["02", "System Design Document", "v1.0",
         "Architecture, component design, integration patterns"],
        ["03", "Security Plan", "v1.0",
         "FIPS 199, OWASP controls, NIST mapping, WAF policy"],
        ["04", "Deployment & Operations Guide", "v1.0",
         "Azure setup, CI/CD, monitoring, runbook procedures"],
        ["05", "User Guide", "v1.0",
         "Field team upload workflow + admin dashboard usage"],
        ["06", "API & Data Reference", "v1.0",
         "REST API endpoints, data model, security headers"],
        ["\u2014", "Executive Summary PPTX", "v2.0",
         "This presentation (14-slide executive briefing)"],
        ["\u2014", "Project Plan XML", "v1.0",
         "MS Project-compatible schedule (10 phases, 90 tasks)"],
    ],
    col_widths=[5, 35, 8, 52],
    font_row=Pt(12),
)
add_footer(slide)


# ── SLIDE 12: RISK ASSESSMENT ───────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_table_slide(slide, "Risk Assessment",
    ["Risk", "Likelihood", "Impact", "Mitigation"],
    [
        ["PIN brute force", "Low", "Medium",
         "Rate limiting + lockout + bcrypt + WAF bot protection"],
        ["Data loss", "Low", "High",
         "Azure automatic backups + blob soft delete + Private Link"],
        ["Network unavailability", "Medium", "Medium",
         "Front Door multi-region routing + health probes"],
        ["Credential exposure", "Low", "High",
         "Key Vault + bcrypt + timing-safe compare + no plaintext"],
        ["CDN cache poisoning", "Low", "Medium",
         "WAF Prevention mode + OWASP DRS 2.1 managed rules"],
        ["DDoS / bot attack", "Medium", "Medium",
         "Front Door WAF + rate limiting + IP restrictions"],
        ["Scale limitations", "Medium", "Low",
         "In-memory rate limit \u2192 Redis migration path ready"],
    ],
    col_widths=[22, 12, 12, 54],
)
add_footer(slide)


# ── SLIDE 13: RECOMMENDATION & APPROVAL ─────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_slide_header(slide, "Recommendation & Approval")

add_title_text(slide,
    "The ASPR Photo Repository application has been successfully deployed "
    "to production. The system meets all functional requirements, adheres "
    "to NIST and OWASP security standards, is protected by Azure Front "
    "Door WAF with OWASP DRS 2.1 ruleset, and operates within full "
    "network isolation via Private Link. The application is recommended "
    "for Authority to Operate (ATO) approval.",
    Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
    font_size=Pt(18), color=WHITE, bold=False)

table_shape = slide.shapes.add_table(
    5, 4, Inches(0.8), Inches(3.4), Inches(11.5), Inches(2.5)
)
table = table_shape.table

headers = ["Role", "Name", "Signature", "Date"]
col_pct = [30, 25, 25, 20]
total = sum(col_pct)
for i, w in enumerate(col_pct):
    table.columns[i].width = int(Inches(11.5) * w / total)

for i, hdr in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE_PRIMARY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True

roles = [
    "Federal Project Sponsor",
    "Information System Security Officer (ISSO)",
    "Authorizing Official (AO)",
    "Technical Lead",
]
for ri, role in enumerate(roles):
    bg = ROW_EVEN if ri % 2 == 0 else ROW_ODD
    for ci in range(4):
        c = table.cell(ri + 1, ci)
        if ci == 0:
            c.text = role
        c.fill.solid()
        c.fill.fore_color.rgb = bg
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE

add_footer(slide)


# ── SLIDE 14: NEXT STEPS ───────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
add_dark_bg(slide)
add_slide_header(slide, "Next Steps")

next_steps = [
    "1.  Complete User Acceptance Testing (UAT) with ASPR field team "
    "representatives",
    "2.  Conduct formal security review and obtain Authority to "
    "Operate (ATO)",
    "3.  Configure Azure Monitor / Application Insights for production "
    "telemetry and alerting",
    "4.  Train operations staff on admin dashboard, PIN management, "
    "and photo workflow",
    "5.  Conduct field pilot during next incident activation or "
    "training exercise",
    "6.  Integrate Login.gov + ID.me external responder authentication "
    "(Phase 2 \u2014 app registration pending)",
    "7.  Plan v1.1 enhancements: interactive map view, offline mode, "
    "batch download improvements",
]

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.9),
                                  Inches(11), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

for i, step in enumerate(next_steps):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

if ASPR_LOGO.exists():
    slide.shapes.add_picture(str(ASPR_LOGO), Inches(0.8), Inches(6.2),
                             height=Inches(0.7))
if LEIDOS_LOGO.exists():
    slide.shapes.add_picture(str(LEIDOS_LOGO), Inches(10.5), Inches(6.3),
                             height=Inches(0.5))

add_footer(slide)


# ══════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
size_kb = OUT.stat().st_size / 1024
print(f"\nExecutive Summary PPTX v2.0 generated: {OUT}")
print(f"Size: {size_kb:.1f} KB")
print(f"Slides: {len(prs.slides)}")
